#!/usr/bin/env python3
"""
Presentation Validation Script

Validates scientific presentations for common issues:
- Slide count vs. duration
- LaTeX compilation
- File size checks
- Basic format validation
"""

import sys
import os
import argparse
import subprocess
from pathlib import Path
from typing import Dict, List, Tuple, Optional

# Try to import PyPDF2 for PDF analysis
try:
    import PyPDF2
    HAS_PYPDF2 = True
except ImportError:
    HAS_PYPDF2 = False

# Try to import python-pptx for PowerPoint analysis
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class PresentationValidator:
    """Validates presentations for common issues."""
    
    # Recommended slide counts by duration (min, recommended, max)
    SLIDE_GUIDELINES = {
        5: (5, 6, 8),
        10: (8, 11, 14),
        15: (13, 16, 20),
        20: (18, 22, 26),
        30: (22, 27, 33),
        45: (32, 40, 50),
        60: (40, 52, 65),
    }
    
    def __init__(self, filepath: str, duration: Optional[int] = None):
        self.filepath = Path(filepath)
        self.duration = duration
        self.file_type = self.filepath.suffix.lower()
        self.issues = []
        self.warnings = []
        self.info = []
        
    def validate(self) -> Dict:
        """Run all validations and return results."""
        print(f"Validating: {self.filepath.name}")
        print(f"File type: {self.file_type}")
        print("=" * 60)
        
        # Check file exists
        if not self.filepath.exists():
            self.issues.append(f"File not found: {self.filepath}")
            return self._format_results()
        
        # File size check
        self._check_file_size()
        
        # Type-specific validation
        if self.file_type == '.pdf':
            self._validate_pdf()
        elif self.file_type in ['.pptx', '.ppt']:
            self._validate_pptx()
        elif self.file_type in ['.tex']:
            self._validate_latex()
        else:
            self.warnings.append(f"Unknown file type: {self.file_type}")
        
        return self._format_results()
    
    def _check_file_size(self):
        """Check if file size is reasonable."""
        size_mb = self.filepath.stat().st_size / (1024 * 1024)
        self.info.append(f"File size: {size_mb:.2f} MB")
        
        if size_mb > 100:
            self.issues.append(
                f"File is very large ({size_mb:.1f} MB). "
                "Consider compressing images."
            )
        elif size_mb > 50:
            self.warnings.append(
                f"File is large ({size_mb:.1f} MB). "
                "May be slow to email or upload."
            )
    
    def _validate_pdf(self):
        """Validate PDF presentation."""
        if not HAS_PYPDF2:
            self.warnings.append(
                "PyPDF2 not installed. Install with: pip install PyPDF2"
            )
            return
        
        try:
            with open(self.filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                num_pages = len(reader.pages)
                
                self.info.append(f"Number of slides: {num_pages}")
                
                # Check slide count against duration
                if self.duration:
                    self._check_slide_count(num_pages)
                
                # Get page size
                first_page = reader.pages[0]
                media_box = first_page.mediabox
                width = float(media_box.width)
                height = float(media_box.height)
                
                # Convert points to inches (72 points = 1 inch)
                width_in = width / 72
                height_in = height / 72
                aspect = width / height
                
                self.info.append(
                    f"Slide dimensions: {width_in:.1f}\" × {height_in:.1f}\" "
                    f"(aspect ratio: {aspect:.2f})"
                )
                
                # Check common aspect ratios
                if abs(aspect - 16/9) < 0.01:
                    self.info.append("Aspect ratio: 16:9 (widescreen)")
                elif abs(aspect - 4/3) < 0.01:
                    self.info.append("Aspect ratio: 4:3 (standard)")
                else:
                    self.warnings.append(
                        f"Unusual aspect ratio: {aspect:.2f}. "
                        "Confirm this matches venue requirements."
                    )
                
        except Exception as e:
            self.issues.append(f"Error reading PDF: {str(e)}")
    
    def _validate_pptx(self):
        """Validate PowerPoint presentation."""
        if not HAS_PPTX:
            self.warnings.append(
                "python-pptx not installed. Install with: pip install python-pptx"
            )
            return
        
        try:
            prs = Presentation(self.filepath)
            num_slides = len(prs.slides)
            
            self.info.append(f"Number of slides: {num_slides}")
            
            # Check slide count against duration
            if self.duration:
                self._check_slide_count(num_slides)
            
            # Get slide dimensions
            width_inches = prs.slide_width / 914400  # EMU to inches
            height_inches = prs.slide_height / 914400
            aspect = prs.slide_width / prs.slide_height
            
            self.info.append(
                f"Slide dimensions: {width_inches:.1f}\" × {height_inches:.1f}\" "
                f"(aspect ratio: {aspect:.2f})"
            )
            
            # Check fonts and text
            self._check_pptx_content(prs)
            
        except Exception as e:
            self.issues.append(f"Error reading PowerPoint: {str(e)}")
    
    def _check_pptx_content(self, prs):
        """Check PowerPoint content for common issues."""
        small_text_slides = []
        many_bullets_slides = []
        
        for idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                text_frame = shape.text_frame
                
                # Check for small fonts
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size and run.font.size.pt < 18:
                            small_text_slides.append(idx)
                            break
                
                # Check for too many bullets
                bullet_count = sum(1 for p in text_frame.paragraphs if p.level == 0)
                if bullet_count > 6:
                    many_bullets_slides.append(idx)
        
        # Report issues
        if small_text_slides:
            unique_slides = sorted(set(small_text_slides))
            self.warnings.append(
                f"Small text (<18pt) found on slides: {unique_slides[:5]}"
                + (" ..." if len(unique_slides) > 5 else "")
            )
        
        if many_bullets_slides:
            unique_slides = sorted(set(many_bullets_slides))
            self.warnings.append(
                f"Many bullets (>6) on slides: {unique_slides[:5]}"
                + (" ..." if len(unique_slides) > 5 else "")
            )
    
    def _validate_latex(self):
        """Validate LaTeX Beamer presentation."""
        self.info.append("LaTeX source file detected")
        
        # Try to compile
        if self._try_compile_latex():
            self.info.append("LaTeX compilation: SUCCESS")
            
            # If PDF was generated, validate it
            pdf_path = self.filepath.with_suffix('.pdf')
            if pdf_path.exists():
                pdf_validator = PresentationValidator(str(pdf_path), self.duration)
                pdf_results = pdf_validator.validate()
                
                # Merge results
                self.info.extend(pdf_results['info'])
                self.warnings.extend(pdf_results['warnings'])
                self.issues.extend(pdf_results['issues'])
        else:
            self.issues.append(
                "LaTeX compilation failed. Check .log file for errors."
            )
    
    def _try_compile_latex(self) -> bool:
        """Try to compile LaTeX file."""
        try:
            # Try pdflatex
            result = subprocess.run(
                ['pdflatex', '-no-shell-escape', '-interaction=nonstopmode', self.filepath.name],
                cwd=self.filepath.parent,
                capture_output=True,
                timeout=60
            )
            return result.returncode == 0
        except (subprocess.TimeoutExpired, FileNotFoundError):
            return False
    
    def _check_slide_count(self, num_slides: int):
        """Check if slide count is appropriate for duration."""
        if self.duration not in self.SLIDE_GUIDELINES:
            # Find nearest duration
            durations = sorted(self.SLIDE_GUIDELINES.keys())
            nearest = min(durations, key=lambda x: abs(x - self.duration))
            min_slides, rec_slides, max_slides = self.SLIDE_GUIDELINES[nearest]
            self.info.append(
                f"Using guidelines for {nearest}-minute talk "
                f"(closest to {self.duration} minutes)"
            )
        else:
            min_slides, rec_slides, max_slides = self.SLIDE_GUIDELINES[self.duration]
        
        self.info.append(
            f"Recommended slides for {self.duration}-minute talk: "
            f"{min_slides}-{max_slides} (optimal: ~{rec_slides})"
        )
        
        if num_slides < min_slides:
            self.warnings.append(
                f"Fewer slides ({num_slides}) than recommended ({min_slides}-{max_slides}). "
                "May have too much time or too little content."
            )
        elif num_slides > max_slides:
            self.warnings.append(
                f"More slides ({num_slides}) than recommended ({min_slides}-{max_slides}). "
                "Likely to run over time."
            )
        else:
            self.info.append(
                f"Slide count ({num_slides}) is within recommended range."
            )
    
    def _format_results(self) -> Dict:
        """Format validation results."""
        return {
            'filepath': str(self.filepath),
            'file_type': self.file_type,
            'info': self.info,
            'warnings': self.warnings,
            'issues': self.issues,
            'valid': len(self.issues) == 0
        }


def print_results(results: Dict):
    """Print validation results in a readable format."""
    print()
    print("=" * 60)
    print("VALIDATION RESULTS")
    print("=" * 60)
    
    # Print info
    if results['info']:
        print("\n📋 Information:")
        for item in results['info']:
            print(f"  • {item}")
    
    # Print warnings
    if results['warnings']:
        print("\n⚠️  Warnings:")
        for item in results['warnings']:
            print(f"  • {item}")
    
    # Print issues
    if results['issues']:
        print("\n❌ Issues:")
        for item in results['issues']:
            print(f"  • {item}")
    
    # Overall status
    print("\n" + "=" * 60)
    if results['valid']:
        print("✅ Validation PASSED")
        if results['warnings']:
            print(f"   ({len(results['warnings'])} warning(s) found)")
    else:
        print("❌ Validation FAILED")
        print(f"   ({len(results['issues'])} issue(s) found)")
    print("=" * 60)


def main():
    parser = argparse.ArgumentParser(
        description='Validate scientific presentations',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s presentation.pdf --duration 15
  %(prog)s slides.pptx --duration 45
  %(prog)s beamer_talk.tex --duration 20

Supported file types:
  - PDF (.pdf)
  - PowerPoint (.pptx, .ppt)
  - LaTeX Beamer (.tex)

Validation checks:
  - Slide count vs. duration
  - File size
  - Slide dimensions
  - Font sizes (PowerPoint)
  - LaTeX compilation (Beamer)
        """
    )
    
    parser.add_argument(
        'filepath',
        help='Path to presentation file (PDF, PPTX, or TEX)'
    )
    
    parser.add_argument(
        '--duration', '-d',
        type=int,
        help='Presentation duration in minutes'
    )
    
    parser.add_argument(
        '--quiet', '-q',
        action='store_true',
        help='Only show issues and warnings'
    )
    
    args = parser.parse_args()
    
    # Validate
    validator = PresentationValidator(args.filepath, args.duration)
    results = validator.validate()
    
    # Print results
    if args.quiet:
        # Only show warnings and issues
        if results['warnings'] or results['issues']:
            print_results(results)
        else:
            print("✅ No issues found")
    else:
        print_results(results)
    
    # Exit with appropriate code
    sys.exit(0 if results['valid'] else 1)


if __name__ == '__main__':
    main()

